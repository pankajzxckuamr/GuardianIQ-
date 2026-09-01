import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Palette
    C_DARK_BG = RGBColor(15, 23, 42)      # #0F172A (Deep Slate / Navy)
    C_LIGHT_BG = RGBColor(248, 250, 252)  # #F8FAFC
    C_CARD_BG = RGBColor(30, 41, 59)      # #1E293B (Card Dark)
    C_LIGHT_CARD = RGBColor(255, 255, 255)# #FFFFFF (Card Light)
    C_ACCENT_BLUE = RGBColor(59, 130, 246)# #3B82F6 (Vibrant Blue)
    C_ACCENT_TEAL = RGBColor(13, 148, 136)# #0D9488 (Teal)
    C_ACCENT_EMERALD = RGBColor(16, 185, 129)# #10B981 (Green)
    C_ACCENT_AMBER = RGBColor(245, 158, 11)# #F59E0B (Amber)
    C_ACCENT_RED = RGBColor(239, 68, 68)  # #EF4444 (Red)
    C_TEXT_WHITE = RGBColor(255, 255, 255)
    C_TEXT_DARK = RGBColor(15, 23, 42)
    C_TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8
    C_TEXT_BODY = RGBColor(51, 65, 85)     # #334155
    C_BORDER_DARK = RGBColor(51, 65, 85)
    C_BORDER_LIGHT = RGBColor(226, 232, 240)

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="GUARDIANIQ ENTERPRISE GOVERNANCE"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(10)
        p_c.font.bold = True
        p_c.font.color.rgb = C_ACCENT_BLUE
        
        # Header title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_DARK

    def set_slide_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1, C_DARK_BG)
    
    # Accent shape
    acc = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.8))
    acc.fill.solid()
    acc.fill.fore_color.rgb = C_ACCENT_BLUE
    acc.line.fill.background()

    # Title & Subtitle box
    t_box = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(3.8))
    tf1 = t_box.text_frame
    tf1.word_wrap = True
    
    p_badge = tf1.paragraphs[0]
    p_badge.text = "PHASE 5 ENGINEERING DELIVERABLE & DEMONSTRATION"
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = C_ACCENT_TEAL
    p_badge.space_after = Pt(10)

    p_main = tf1.add_paragraph()
    p_main.text = "GuardianIQ Policy & Runtime ENFORCE Layer"
    p_main.font.size = Pt(36)
    p_main.font.bold = True
    p_main.font.color.rgb = C_TEXT_WHITE
    p_main.space_after = Pt(10)

    p_sub = tf1.add_paragraph()
    p_sub.text = "Deterministic Multi-Layer Policy Hierarchy, Runtime Enforcement Gateway, Cryptographic TOCTOU Protection & Non-Authoritative Simulator"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = C_TEXT_MUTED
    p_sub.space_after = Pt(24)

    p_meta = tf1.add_paragraph()
    p_meta.text = "Presenter: Aayush Kumar (Fullstack & Systems Lead)  |  Engineering Team: Pankaj, Jitendra  |  Version: 1.0 (Production Verified)"
    p_meta.font.size = Pt(11)
    p_meta.font.color.rgb = RGBColor(203, 213, 225)

    # ==========================================
    # SLIDE 2: Executive Overview & The Problem Solved
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2, C_LIGHT_BG)
    add_header(s2, "Executive Overview: Solving The AI Governance Gap")

    cards_s2 = [
        ("The Challenge in Autonomous AI", 
         "• Unbounded Agent Autonomy: LLM agents calling tools and APIs without deterministic controls.\n"
         "• Sensitive Data Leakage: Raw PII/PHI sent to third-party model providers.\n"
         "• Time-of-Check Vulnerabilities (TOCTOU): Approved actions modified before target execution.\n"
         "• Disjointed Audit Lineage: Missing correlation across request, approval, and execution.",
         C_ACCENT_RED, Inches(0.8)),
        ("The GuardianIQ Solution (Phase 5)", 
         "• Deterministic Policy Engine: Safe AST condition evaluation without dynamic code execution.\n"
         "• Hard Boundary Guards: Autonomy limits, tool caps, data masking, and model restrictions.\n"
         "• Single-Use Cryptographic Tokens: SHA-256 context hashing preventing replay and tampering.\n"
         "• Transactional Outbox Events: Complete audit trail sharing unified correlation_id.",
         C_ACCENT_BLUE, Inches(4.8)),
        ("Key Business & Technical Outcomes", 
         "• 100% Fail-Closed Security: Unknown states, syntax errors, and timeouts default to DENY.\n"
         "• Enterprise Compliance: HIPAA, GDPR, SOC2, and statutory 7-year audit retention ready.\n"
         "• Zero Performance Degradation: In-memory caching with sub-millisecond policy checks.\n"
         "• 82/82 Automated Tests Passing: Zero regressions across all master deliverables.",
         C_ACCENT_EMERALD, Inches(8.8))
    ]

    for title, text, color, left in cards_s2:
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.7), Inches(5.3))
        card.fill.solid()
        card.fill.fore_color.rgb = C_LIGHT_CARD
        card.line.color.rgb = C_BORDER_LIGHT
        card.line.width = Pt(1.5)
        
        # Top color accent bar
        bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), Inches(1.7), Inches(3.3), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        tb = s2.shapes.add_textbox(left + Inches(0.2), Inches(1.9), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_DARK
        p_t.space_after = Pt(10)
        
        p_b = tf.add_paragraph()
        p_b.text = text
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = C_TEXT_BODY
        p_b.line_spacing = 1.25

    # ==========================================
    # SLIDE 3: System Topology & Runtime Architecture
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3, C_LIGHT_BG)
    add_header(s3, "Phase 5 Architecture: Runtime Enforcement Pipeline")

    stages = [
        ("1. Client / Agent Call", "Agent invokes action or tool via /api/v1/enforce/execute. Direct target calls are blocked.", Inches(0.8), "0F172A"),
        ("2. Boundary Guards", "Evaluates Autonomy Level, Kill Switch, Tool Limits, Data Classification & Model Provider.", Inches(3.25), "1E3A8A"),
        ("3. Policy Resolver", "Traverses Direct -> Department -> Global hierarchy; evaluates AST rules; combines decisions.", Inches(5.7), "047857"),
        ("4. TOCTOU & AuthZ", "Generates SHA-256 hashes; validates context equality; issues single-use token.", Inches(8.15), "B45309"),
        ("5. Target & Outbox", "Executes target adapter; publishes immutable audit event to transactional outbox.", Inches(10.6), "4338CA")
    ]

    for title, desc, left, col_hex in stages:
        c_rgb = RGBColor(int(col_hex[:2], 16), int(col_hex[2:4], 16), int(col_hex[4:], 16))
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(2.2), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = C_LIGHT_CARD
        card.line.color.rgb = c_rgb
        card.line.width = Pt(2)
        
        tb = s3.shapes.add_textbox(left + Inches(0.1), Inches(1.8), Inches(2.0), Inches(3.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = c_rgb
        p_t.space_after = Pt(8)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = C_TEXT_BODY
        p_d.line_spacing = 1.2

    # Bottom summary callout card
    b_card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.4))
    b_card.fill.solid()
    b_card.fill.fore_color.rgb = RGBColor(241, 245, 249)
    b_card.line.color.rgb = C_BORDER_LIGHT
    
    b_tb = s3.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.2))
    b_tf = b_tb.text_frame
    b_tf.word_wrap = True
    p_bt = b_tf.paragraphs[0]
    p_bt.text = "CORE ARCHITECTURAL GUARANTEES"
    p_bt.font.size = Pt(11)
    p_bt.font.bold = True
    p_bt.font.color.rgb = C_ACCENT_BLUE
    p_bt.space_after = Pt(4)
    
    p_bd = b_tf.add_paragraph()
    p_bd.text = "• Fail-Closed: Engine timeouts, network breaks, and AST parsing errors default to DENY.\n" \
                "• Strict Precedence: DENY > REQUIRE_APPROVAL > ESCALATE > ALLOW_WITH_OBLIGATIONS > ALLOW.\n" \
                "• Unified Correlation: Every step in the execution pipeline shares the exact same correlation_id for zero-loss audit forensics."
    p_bd.font.size = Pt(9.5)
    p_bd.font.color.rgb = C_TEXT_BODY

    # ==========================================
    # SLIDE 4: Policy Hierarchy, Versioning & In-UI Rule Builder
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4, C_LIGHT_BG)
    add_header(s4, "Multi-Layer Policy Hierarchy & Rule Evaluator Engine")

    cards_s4 = [
        ("Hierarchical Resolution Order",
         "1. DIRECT Bindings (Entity-Specific)\n"
         "   Bound to specific Agent, Tool, Model, or Data Source (Highest priority).\n\n"
         "2. DEPARTMENT Bindings\n"
         "   Inherited by all assets belonging to the department.\n\n"
         "3. GLOBAL (*) Bindings\n"
         "   System-wide baseline policies applying to every asset.\n\n"
         "★ Mandatory Policy Lock: Mandatory parent policies cannot be overridden or relaxed by lower-level bindings.",
         Inches(0.8), Inches(5.6)),
        ("Deterministic AST Rule Evaluator",
         "• Safe AST Condition Parsing:\n"
         "  Expressions like `agent.risk_level == 'CRITICAL' and tool.access_mode == 'WRITE'` are parsed into syntax trees.\n\n"
         "• Zero Dynamic Execution:\n"
         "  Python `eval()` and `exec()` are strictly banned.\n\n"
         "• In-UI Rule Builder & Versioning:\n"
         "  Supports DRAFT -> ACTIVE -> DEPRECATED -> ARCHIVED lifecycle.\n"
         "  Active versions are immutable to guarantee historic reproducibility.",
         Inches(6.7), Inches(5.8))
    ]

    for title, text, left, width in cards_s4:
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), width, Inches(5.3))
        card.fill.solid()
        card.fill.fore_color.rgb = C_LIGHT_CARD
        card.line.color.rgb = C_BORDER_LIGHT
        card.line.width = Pt(1.5)
        
        tb = s4.shapes.add_textbox(left + Inches(0.25), Inches(1.7), width - Inches(0.5), Inches(4.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = C_ACCENT_BLUE
        p_t.space_after = Pt(10)
        
        p_b = tf.add_paragraph()
        p_b.text = text
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = C_TEXT_BODY
        p_b.line_spacing = 1.25

    # ==========================================
    # SLIDE 5: Agent Boundary, Tool, Data & Model Guards
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5, C_LIGHT_BG)
    add_header(s5, "Hard Security Boundaries: Agent, Tool, Data & Model Guards")

    guards = [
        ("Agent Autonomy Guard", 
         "• Autonomy Ceilings:\n  RECOMMEND_ONLY (no write)\n  HUMAN_IN_THE_LOOP (approval)\n  AUTONOMOUS_WITH_BOUNDS\n\n• Instant Kill Switch:\n  One-click administrative toggle instantly blocks all actions.",
         Inches(0.8), "EF4444"),
        ("Tool Permission Guard", 
         "• Relationship Validation:\n  Verifies explicit USES_TOOL binding.\n\n• Capability Tags & Ceilings:\n  READ / WRITE / ADMIN modes.\n  Parameter thresholds (e.g. max refund amount <= $10,000).",
         Inches(3.75), "3B82F6"),
        ("Data Permission Guard", 
         "• Classification Hierarchy:\n  PUBLIC, INTERNAL, CONFIDENTIAL, RESTRICTED.\n\n• Field-Level Transforms:\n  Applies dynamic MASK, REDACT, TOKENIZE, or HASH to sensitive fields before exposure.",
         Inches(6.7), "10B981"),
        ("Model Provider Guard", 
         "• Provider Whitelist:\n  Prevents sending proprietary data to unapproved public LLMs.\n\n• Fallback Trap:\n  Blocks rogue fallback calls to unauthorized fallback models.",
         Inches(9.65), "8B5CF6")
    ]

    for title, desc, left, col_hex in guards:
        c_rgb = RGBColor(int(col_hex[:2], 16), int(col_hex[2:4], 16), int(col_hex[4:], 16))
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(2.8), Inches(5.3))
        card.fill.solid()
        card.fill.fore_color.rgb = C_LIGHT_CARD
        card.line.color.rgb = C_BORDER_LIGHT
        card.line.width = Pt(1.5)
        
        bar = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), Inches(1.7), Inches(2.5), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = c_rgb
        bar.line.fill.background()

        tb = s5.shapes.add_textbox(left + Inches(0.15), Inches(1.9), Inches(2.5), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_DARK
        p_t.space_after = Pt(10)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = C_TEXT_BODY
        p_d.line_spacing = 1.25

    # ==========================================
    # SLIDE 6: Cryptographic TOCTOU Protection & Single-Use Tokens
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6, C_LIGHT_BG)
    add_header(s6, "Security Hardening: Cryptographic TOCTOU & Replay Defense")

    s6_cards = [
        ("The Time-of-Check to Time-of-Use (TOCTOU) Threat",
         "In multi-agent systems, a critical security vulnerability occurs when an action is evaluated or approved at Time T0 (e.g. refund $500), but modified at Time T1 right before execution (e.g. refund $50,000).\n\n"
         "Without cryptographic request-binding, traditional governance systems are blind to payload tampering between approval and execution.",
         Inches(0.8), Inches(5.6)),
        ("GuardianIQ Multi-Hash Authorization Service",
         "1. Tri-Hash Request Fingerprint:\n"
         "   • context_hash = SHA-256(request_payload)\n"
         "   • relationship_hash = SHA-256(active_graph)\n"
         "   • policy_hash = SHA-256(active_rules)\n\n"
         "2. Single-Use Runtime Authorization Token:\n"
         "   • Issues a cryptographic token bound to the exact Tri-Hash fingerprint with a 60-second TTL.\n\n"
         "3. Re-Verification & Replay Blocker:\n"
         "   • Execution engine re-hashes the live payload immediately before target invocation.\n"
         "   • Any hash mismatch rejects execution. Tokens are marked consumed_at to block replay attacks.",
         Inches(6.7), Inches(5.8))
    ]

    for title, text, left, width in s6_cards:
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), width, Inches(5.3))
        card.fill.solid()
        card.fill.fore_color.rgb = C_LIGHT_CARD
        card.line.color.rgb = C_BORDER_LIGHT
        card.line.width = Pt(1.5)
        
        tb = s6.shapes.add_textbox(left + Inches(0.25), Inches(1.7), width - Inches(0.5), Inches(4.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = C_ACCENT_BLUE
        p_t.space_after = Pt(10)
        
        p_b = tf.add_paragraph()
        p_b.text = text
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = C_TEXT_BODY
        p_b.line_spacing = 1.25

    # ==========================================
    # SLIDE 7: Non-Authoritative Simulator & Decision Trace
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7, C_LIGHT_BG)
    add_header(s7, "Developer & Operator Experience: Enforcement Simulator")

    sim_features = [
        ("Zero Side-Effects", "Simulations test complex multi-agent payloads without mutating live state or invoking external target adapters.", Inches(0.8)),
        ("Multi-Layer Decision Trace", "Visualizes step-by-step resolution across Relationship Guards, Hard Boundaries, and Policy Rules.", Inches(3.75)),
        ("Obligation Inspections", "Inspects required downstream actions like field masking, telemetry tags, or supervisor escalations.", Inches(6.7)),
        ("Remediation Guidance", "Returns exact AST rule codes and failure reasons so developers can fix permission mismatches fast.", Inches(9.65))
    ]

    for title, desc, left in sim_features:
        card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(2.8), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = C_LIGHT_CARD
        card.line.color.rgb = C_BORDER_LIGHT
        card.line.width = Pt(1.5)
        
        tb = s7.shapes.add_textbox(left + Inches(0.15), Inches(1.65), Inches(2.5), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = C_ACCENT_TEAL
        p_t.space_after = Pt(6)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = C_TEXT_BODY

    # Simulator UI Walkthrough Box
    sim_ui_box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.7))
    sim_ui_box.fill.solid()
    sim_ui_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    sim_tb = s7.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.3), Inches(2.5))
    sim_tf = sim_tb.text_frame
    sim_tf.word_wrap = True
    
    p_st = sim_tf.paragraphs[0]
    p_st.text = "LIVE SIMULATION WORKFLOW IN THE DEMO"
    p_st.font.size = Pt(12)
    p_st.font.bold = True
    p_st.font.color.rgb = C_ACCENT_BLUE
    p_st.space_after = Pt(6)
    
    p_sd = sim_tf.add_paragraph()
    p_sd.text = "1. Select Agent & Tool / Data Source from autocomplete dropdowns.\n" \
                "2. Paste JSON action payload (e.g. {\"amount\": 15000, \"currency\": \"USD\", \"reason\": \"customer_satisfaction\"}).\n" \
                "3. Click 'Run Enforcement Simulation' -> Sub-50ms execution trace returns matched rules, decision badges, and obligations.\n" \
                "4. Live Decision Breakdown: Green (ALLOW), Yellow (REQUIRE_APPROVAL), Red (DENY) with rule codes and AST condition diffs."
    p_sd.font.size = Pt(10)
    p_sd.font.color.rgb = RGBColor(226, 232, 240)
    p_sd.line_spacing = 1.3

    # ==========================================
    # SLIDE 8: Automated Test Verification Evidence
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8, C_LIGHT_BG)
    add_header(s8, "Verification Evidence: 82/82 Automated Tests Passing")

    test_groups = [
        ("Database & Migrations", "6 Tests", "PostgreSQL schema, foreign keys, triggers, Alembic forward & rollback.", Inches(0.8)),
        ("Binding Hierarchy", "18 Tests", "Direct, Department, Global resolution, precedence order, mandatory lock.", Inches(3.75)),
        ("Boundary & Guards", "19 Tests", "Autonomy, Kill Switch, Tool ceilings, Data classification, Masking.", Inches(6.7)),
        ("Gateway & TOCTOU", "11 Tests", "Exact-request authorization, SHA-256 hashes, single-use token replay.", Inches(9.65)),
        ("Approvals & Events", "12 Tests", "Workflow approval bridge, exception lookups, Outbox event correlation.", Inches(0.8)),
        ("Simulator & API", "11 Tests", "Non-authoritative simulation, REST validation, OpenAPI contract checks.", Inches(3.75)),
        ("Security & Isolation", "5 Tests", "Cross-tenant security, agent identity spoofing, secret stripping.", Inches(6.7)),
        ("Overall Scorecard", "100% Pass", "82 Passed, 0 Failed, 0 Regressions across all Phase 5 modules.", Inches(9.65))
    ]

    for idx, (title, count, desc, left) in enumerate(test_groups):
        top = Inches(1.5) if idx < 4 else Inches(4.2)
        card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = C_LIGHT_CARD
        card.line.color.rgb = C_BORDER_LIGHT
        card.line.width = Pt(1.5)
        
        tb = s8.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), Inches(2.5), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_DARK
        
        p_c = tf.add_paragraph()
        p_c.text = count
        p_c.font.size = Pt(16)
        p_c.font.bold = True
        p_c.font.color.rgb = C_ACCENT_EMERALD if "Pass" in count or "82" in count else C_ACCENT_BLUE
        p_c.space_after = Pt(4)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = C_TEXT_BODY

    # ==========================================
    # SLIDE 9: Live Demo Roadmap & Key Takeaways
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9, C_DARK_BG)
    
    # Title
    t_box9 = s9.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
    tf9 = t_box9.text_frame
    p_t9 = tf9.paragraphs[0]
    p_t9.text = "Tomorrow's Live Demonstration Roadmap"
    p_t9.font.size = Pt(24)
    p_t9.font.bold = True
    p_t9.font.color.rgb = C_TEXT_WHITE

    demo_steps = [
        ("Step 1: Policy Management", "Navigate to /policies. Review 10-per-page pagination, target filtering, and activate version v1 of High-Risk Agent Controls.", Inches(0.8), Inches(1.6)),
        ("Step 2: In-UI Rule Builder", "Create version v2 with new rate limit rule. Create new data policy POL-FIN-REFUND-001 with initial rules.", Inches(0.8), Inches(2.9)),
        ("Step 3: Attach & Inspect", "Attach policy to Autonomous Refund Agent. Verify Direct -> Dept -> Global hierarchy in Applicable Policies Inspector.", Inches(0.8), Inches(4.2)),
        ("Step 4: Agent Boundary Tab", "Navigate to Agent Detail. Toggle Kill Switch, inspect tool ceilings ($10k limit), and review data field permissions.", Inches(6.8), Inches(1.6)),
        ("Step 5: Live Simulator Run", "Run 4 simulation scenarios on /enforcement-simulation: Permitted ($2.5k), High-Value Approval ($15k), Masked PII, and Denied.", Inches(6.8), Inches(2.9)),
        ("Step 6: Audit Forensics", "Inspect live enforcement event trail on Agent Enforcement History and verify unified correlation_id lineage.", Inches(6.8), Inches(4.2))
    ]

    for title, desc, left, top in demo_steps:
        card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER_DARK
        
        tb = s9.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), Inches(5.3), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = C_ACCENT_TEAL
        p_t.space_after = Pt(2)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = RGBColor(203, 213, 225)

    # Bottom summary
    bot_box = s9.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.8))
    bot_tf = bot_box.text_frame
    p_b = bot_tf.paragraphs[0]
    p_b.text = "★ Phase 5 is Production-Ready, Fail-Closed, Fully Tested, and Certified for Live Demonstration."
    p_b.font.size = Pt(13)
    p_b.font.bold = True
    p_b.font.color.rgb = C_ACCENT_EMERALD
    p_b.alignment = PP_ALIGN.CENTER

    output_path = r"c:\Users\aayus\Desktop\GuardianIQ--1\docs\Phase 5\GuardianIQ_Phase5_Executive_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
