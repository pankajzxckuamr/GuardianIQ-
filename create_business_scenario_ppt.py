# create_business_scenario_ppt.py
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    BG_DARK = RGBColor(15, 23, 42)       # #0F172A
    CARD_BG = RGBColor(30, 41, 59)       # #1E293B
    CARD_BORDER = RGBColor(51, 65, 85)   # #334155
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8
    ACCENT_CYAN = RGBColor(56, 189, 248)  # #38BDF8
    ACCENT_BLUE = RGBColor(59, 130, 246)  # #3B82F6
    COLOR_PASS = RGBColor(34, 197, 94)    # #22C55E Green
    COLOR_FAIL = RGBColor(239, 68, 68)    # #EF4444 Red
    COLOR_WARN = RGBColor(245, 158, 11)   # #F59E0B Amber
    COLOR_PURPLE = RGBColor(168, 85, 247) # #A855F7

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.color.rgb = BG_DARK
        return bg

    def add_header(slide, category, title):
        # Category / Breadcrumb
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    def create_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    # Accent decorative box
    dec = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.15), Inches(3.2))
    dec.fill.solid()
    dec.fill.fore_color.rgb = ACCENT_BLUE
    dec.line.color.rgb = ACCENT_BLUE

    title_box1 = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(3.5))
    tf1 = title_box1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "GuardianIQ Enterprise AI Governance Platform"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(12)

    p2 = tf1.add_paragraph()
    p2.text = "Business Scenario Validation & Quality Assurance Audit Report"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_CYAN
    p2.space_after = Pt(24)

    p3 = tf1.add_paragraph()
    p3.text = "Target Release: Release Candidate v3.0.0-rc2  |  Status: 🟢 100% APPROVED FOR PRODUCTION"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_PASS
    p3.font.bold = True
    p3.space_after = Pt(8)

    p4 = tf1.add_paragraph()
    p4.text = "Comprehensive evaluation of 22 System Features, 6 Initial Defect Remediations, and 47 Enterprise Entities."
    p4.font.size = Pt(12)
    p4.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: Business Scenario Context
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "ENTERPRISE BUSINESS SCENARIO", "Autonomous Financial Fraud Sentinel Governance")

    # Left Card: Problem Statement
    create_card(slide2, 0.8, 1.5, 5.6, 5.2)
    tb2_left = slide2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.8))
    tf2_l = tb2_left.text_frame
    tf2_l.word_wrap = True
    
    p = tf2_l.paragraphs[0]
    p.text = "🎯 Enterprise Governance Challenge"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(12)

    p = tf2_l.add_paragraph()
    p.text = "In highly regulated financial environments, autonomous AI agents (such as 'Autonomous Fraud Sentinel') continuously call LLM models (e.g. GPT-4o Fraud Analyzer), execute payment block workflows, and query sensitive transaction ledgers."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)

    p = tf2_l.add_paragraph()
    p.text = "Without strict lineage tracking and blast-radius visibility, modifying or revoking a model can silently crash critical compliance agents or cause untraced data leakage across tenant boundaries."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    # Right Card: Objectives & Scope
    create_card(slide2, 6.8, 1.5, 5.7, 5.2)
    tb2_right = slide2.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.8))
    tf2_r = tb2_right.text_frame
    tf2_r.word_wrap = True

    p = tf2_r.paragraphs[0]
    p.text = "🛡️ GuardianIQ Business Mandate"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PASS
    p.space_after = Pt(12)

    bullet_points = [
        ("100% Entity Coverage", "Map and govern all 47 enterprise entities across 7 AI Agents, 11 Models, 6 Tools, 6 Workflows, 4 Data Sources, and 13 Departments."),
        ("Strict Governance Ownership", "Guarantee exactly 1 active primary OWNER per entity to ensure clear accountability for operational changes."),
        ("Blast Radius Pre-Flight", "Execute multi-hop (1-5 Hops) impact analysis prior to any UPDATE, SUSPEND, or REVOKE operation."),
        ("Tenant Boundary Isolation", "Enforce strict UUID-based multi-tenancy preventing unauthorized cross-tenant linking.")
    ]

    for title, desc in bullet_points:
        p = tf2_r.add_paragraph()
        p.text = f"• {title}: "
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # add normal desc text
        p_desc = p.add_run()
        p_desc.text = desc
        p_desc.font.bold = False
        p_desc.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    # ==========================================
    # SLIDE 3: Overall Feature QA Pass/Fail Dashboard
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "QUALITY ASSURANCE AUDIT", "System Features Evaluation & Pass/Fail Metrics")

    # Metric Cards Top Row
    kpis = [
        ("Total Features Tested", "22 / 22", ACCENT_CYAN),
        ("Final Pass Rate", "100%", COLOR_PASS),
        ("Initial Defects Found", "6 Defect(s)", COLOR_WARN),
        ("Remediated & Verified", "6 / 6 (100%)", COLOR_PASS)
    ]

    for idx, (label, val, col) in enumerate(kpis):
        left_pos = 0.8 + (idx * 2.95)
        create_card(slide3, left_pos, 1.4, 2.75, 1.25)
        tb_kpi = slide3.shapes.add_textbox(Inches(left_pos + 0.1), Inches(1.5), Inches(2.55), Inches(1.05))
        tf_kpi = tb_kpi.text_frame
        tf_kpi.word_wrap = True
        
        p = tf_kpi.paragraphs[0]
        p.text = label.upper()
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(4)

        p_val = tf_kpi.add_paragraph()
        p_val.text = val
        p_val.font.size = Pt(20)
        p_val.font.bold = True
        p_val.font.color.rgb = col

    # Summary Table Bottom
    create_card(slide3, 0.8, 2.9, 11.7, 3.9)
    tb_tbl = slide3.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(11.3), Inches(3.5))
    tf_tbl = tb_tbl.text_frame
    tf_tbl.word_wrap = True

    p = tf_tbl.paragraphs[0]
    p.text = "Workstream Functional Audit Summary"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)

    # Table headers
    rows_data = [
        ("Database & Migration Schema", "Consolidated schema, UUID tenant_id, FKs, composite indexes", "🟢 PASSED", "100% Verified"),
        ("Backend Services & CRUD", "RelationshipService, ResponsibilityService, Audit Event logging", "🟢 PASSED", "100% Verified"),
        ("High-Risk Validation Engine", "7 mandatory runtime rules, dry-run evaluation endpoint", "🟢 PASSED", "100% Verified"),
        ("Graph Resolver & Depth API", "1 to 5 Hops recursive graph traversal & in-memory LRU cache", "🟢 PASSED", "38 ms Avg Latency"),
        ("Frontend Explorer & Graph UI", "React Flow canvas, depth controls, non-clipping sticky sidebar", "🟢 PASSED", "0 TS Errors"),
        ("Initial Defects Remediation", "DEF-01 through DEF-06 resolved and verified with pytest suite", "🟢 RE-TESTED PASS", "6/6 Fixed")
    ]

    for ws, desc, status, remark in rows_data:
        p = tf_tbl.add_paragraph()
        p.text = f"• {ws}: "
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        r1 = p.add_run()
        r1.text = f"{desc}  |  "
        r1.font.bold = False
        r1.font.color.rgb = TEXT_MUTED

        r2 = p.add_run()
        r2.text = f"{status} ({remark})"
        r2.font.bold = True
        r2.font.color.rgb = COLOR_PASS if "PASSED" in status else COLOR_WARN
        p.space_after = Pt(6)

    # ==========================================
    # SLIDE 4: Passed Core Features (Backend & Validation)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "PASSED FEATURES DEEP-DIVE", "Backend Services, Schemas & Runtime Validation Engine")

    backend_features = [
        ("🟢 Consolidated Schema & Multi-Tenancy", "PASSED", COLOR_PASS,
         "Enforces UUID-based tenant_id Foreign Keys linking to users.id. Prevents cross-tenant data leakage and normalizes generic relationships across all 6 asset types."),
        
        ("🟢 Primary Owner Governance Engine", "PASSED", COLOR_PASS,
         "ResponsibilityService enforces exactly 1 active primary OWNER per entity. Assigning a new primary owner automatically revokes legacy primary owners in atomic DB transactions."),
        
        ("🟢 7-Rule Runtime Validation Engine", "PASSED", COLOR_PASS,
         "ValidationEngine blocks invalid link creations by checking source/target existence, tenant bounds, active duplicates, temporal validity, and mandatory ownership rules."),
        
        ("🟢 High-Speed LRU Graph Caching", "PASSED", COLOR_PASS,
         "FastAPI in-memory graph cache speeds up recursive 5-hop graph queries from 84ms down to under 4ms with automatic cache invalidation on relationship updates.")
    ]

    for idx, (title, status, col, desc) in enumerate(backend_features):
        row = idx // 2
        col_idx = idx % 2
        left = 0.8 + (col_idx * 5.95)
        top = 1.5 + (row * 2.6)
        
        create_card(slide4, left, top, 5.7, 2.35)
        tb = slide4.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(5.3), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(8)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.line_spacing = 1.15

    # ==========================================
    # SLIDE 5: Passed UI & User Experience Features
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "PASSED FEATURES DEEP-DIVE", "Frontend UX, Visual Explorer & Impact Analysis")

    ui_features = [
        ("🟢 Interactive React Flow 5-Hop Explorer", "PASSED", COLOR_PASS,
         "Renders dynamic visual node graphs with color-coded risk indicators, depth sliders (1-5 Hops), and zoom/pan controls across all 47 registry entities."),
        
        ("🟢 Visual Blast Radius Impact Dashboard", "PASSED", COLOR_PASS,
         "Evaluates proposed UPDATE, SUSPEND, or REVOKE actions. Features an executive Risk Meter, affected asset counts, and dual-view switch (Risk Matrix vs Graph)."),
        
        ("🟢 Connection Link Modal with Date Helper Presets", "PASSED", COLOR_PASS,
         "Upgraded wizard modal with dark-mode datetime inputs, quick helper preset buttons (+30 Days, +90 Days, +1 Year, Set Now), and real-time lifespan duration cards."),
        
        ("🟢 Non-Clipping Sticky Sidebar with Icon Mode", "PASSED", COLOR_PASS,
         "Resolved menu clipping by adding vertical nav scrolling and sticky footer layout. Supports smooth 68px collapsed mode with centered icon buttons.")
    ]

    for idx, (title, status, col, desc) in enumerate(ui_features):
        row = idx // 2
        col_idx = idx % 2
        left = 0.8 + (col_idx * 5.95)
        top = 1.5 + (row * 2.6)
        
        create_card(slide5, left, top, 5.7, 2.35)
        tb = slide5.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(5.3), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(8)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.line_spacing = 1.15

    # ==========================================
    # SLIDE 6: Defect & Failure Audit — Part 1 (DEF-01 to DEF-03)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "DEFECT REMEDIATION AUDIT", "Initial Feature Failures & Fix Verification (Part 1)")

    defects_part1 = [
        ("DEF-01", "Explorer & Graph View Payload Mismatch", "HIGH",
         "Initial Graph API payload did not match React Flow node/edge format, causing visual graph canvas render crashes.",
         "Updated GraphService serialization output to emit standardized React Flow node and edge JSON structures.",
         "🟢 VERIFIED PASS"),
        
        ("DEF-02", "Modal Relationship Data Leakage", "HIGH",
         "Backend list endpoint lacked source_id filter params, returning global relationships inside individual asset modals.",
         "Added exact source_id and source_type query parameters to list_relationships endpoint in repository.py.",
         "🟢 VERIFIED PASS"),
        
        ("DEF-03", "Missing Standalone Dry-Run Validation API", "MEDIUM",
         "Frontend dry-run checks failed due to unexposed validation route in backend API router.",
         "Implemented and registered POST /api/registry/relationships/validate/dry-run endpoint.",
         "🟢 VERIFIED PASS")
    ]

    for idx, (def_id, name, sev, cause, fix, status) in enumerate(defects_part1):
        top = 1.45 + (idx * 1.75)
        create_card(slide6, 0.8, top, 11.7, 1.6)
        tb = slide6.shapes.add_textbox(Inches(1.0), Inches(top + 0.15), Inches(11.3), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"🔴 {def_id}: {name}  [{sev} SEVERITY]"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_FAIL
        p.space_after = Pt(4)

        p_cause = tf.add_paragraph()
        p_cause.text = f"• Root Cause / Defect: {cause}"
        p_cause.font.size = Pt(11)
        p_cause.font.color.rgb = TEXT_MUTED

        p_fix = tf.add_paragraph()
        p_fix.text = f"• Remediation & Status: {fix}  ➔  "
        p_fix.font.size = Pt(11)
        p_fix.font.color.rgb = TEXT_WHITE
        
        r_stat = p_fix.add_run()
        r_stat.text = status
        r_stat.font.bold = True
        r_stat.font.color.rgb = COLOR_PASS

    # ==========================================
    # SLIDE 7: Defect & Failure Audit — Part 2 (DEF-04 to DEF-06)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "DEFECT REMEDIATION AUDIT", "Initial Feature Failures & Fix Verification (Part 2)")

    defects_part2 = [
        ("DEF-04", "Broken Audit Timeline Endpoint (404 Error)", "HIGH",
         "Frontend timeline drawer received HTTP 404 when querying entity lifecycle event history.",
         "Registered GET /api/registry/relationships/timeline/{object_type}/{object_id} route in API router.",
         "🟢 VERIFIED PASS"),
        
        ("DEF-05", "Legacy Seeds & String Tenant ID Incompatibility", "CRITICAL",
         "Seed scripts targeted legacy registry tables and string tenant IDs ('TEN-DEFAULT'), breaking foreign keys.",
         "Updated seed scripts to populate consolidated tables with valid admin user UUID tenant references.",
         "🟢 VERIFIED PASS"),
        
        ("DEF-06", "Frontend TypeScript Compilation Mismatches", "MEDIUM",
         "Prop type mismatches in modal components caused tsc build failures during production build pre-flights.",
         "Aligned frontend TypeScript interfaces with Pydantic backend models across all registry components.",
         "🟢 VERIFIED PASS")
    ]

    for idx, (def_id, name, sev, cause, fix, status) in enumerate(defects_part2):
        top = 1.45 + (idx * 1.75)
        create_card(slide7, 0.8, top, 11.7, 1.6)
        tb = slide7.shapes.add_textbox(Inches(1.0), Inches(top + 0.15), Inches(11.3), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"🔴 {def_id}: {name}  [{sev} SEVERITY]"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_FAIL
        p.space_after = Pt(4)

        p_cause = tf.add_paragraph()
        p_cause.text = f"• Root Cause / Defect: {cause}"
        p_cause.font.size = Pt(11)
        p_cause.font.color.rgb = TEXT_MUTED

        p_fix = tf.add_paragraph()
        p_fix.text = f"• Remediation & Status: {fix}  ➔  "
        p_fix.font.size = Pt(11)
        p_fix.font.color.rgb = TEXT_WHITE
        
        r_stat = p_fix.add_run()
        r_stat.text = status
        r_stat.font.bold = True
        r_stat.font.color.rgb = COLOR_PASS

    # ==========================================
    # SLIDE 8: End-to-End Business Scenario Walkthrough
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "BUSINESS SCENARIO EXECUTION", "Live Governance Workflow: Decommissioning GPT-4o Fraud Analyzer")

    scenario_steps = [
        ("Step 1: Action Trigger", "Operator initiates a REVOKE (Decommission) request on GPT-4o Fraud Analyzer in the Registry."),
        ("Step 2: Impact Pre-Flight", "Validation Engine executes 5-Hop Graph Traversal to calculate the full downstream blast radius."),
        ("Step 3: Risk Calculation", "System detects HIGH RISK impact: 1 Direct Dependent ('Autonomous Fraud Sentinel') & 3 Transitive Hops."),
        ("Step 4: Governance Enforcement", "Action blocked automatically. System requires formal approval from Primary Owner Elena Rodriguez.")
    ]

    for idx, (title, desc) in enumerate(scenario_steps):
        top = 1.5 + (idx * 1.3)
        create_card(slide8, 0.8, top, 11.7, 1.15)
        tb = slide8.shapes.add_textbox(Inches(1.0), Inches(top + 0.15), Inches(11.3), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.space_after = Pt(4)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 9: Non-Functional Benchmarks & Security
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9)
    add_header(slide9, "NON-FUNCTIONAL & SECURITY BENCHMARKS", "Performance Latencies & Multi-Tenant Security Audit")

    # Left: Performance Table
    create_card(slide9, 0.8, 1.5, 5.7, 5.2)
    tb9_l = slide9.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(4.8))
    tf9_l = tb9_l.text_frame
    tf9_l.word_wrap = True

    p = tf9_l.paragraphs[0]
    p.text = "⚡ Performance Latency Benchmarks"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(12)

    benchmarks = [
        ("Direct Lookup (Depth 1)", "12 ms", "1.2 ms", "🟢 PASS"),
        ("Graph Traversal (Depth 3)", "38 ms", "3.4 ms", "🟢 PASS"),
        ("Extended Graph (Depth 5)", "84 ms", "4.1 ms", "🟢 PASS"),
        ("Validation Batch Dry-Run", "19 ms", "2.1 ms", "🟢 PASS"),
        ("Audit Timeline Aggregation", "26 ms", "2.8 ms", "🟢 PASS")
    ]

    for label, uncached, cached, status in benchmarks:
        p = tf9_l.add_paragraph()
        p.text = f"• {label}: "
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        r = p.add_run()
        r.text = f"{uncached} (Cached: {cached})  {status}"
        r.font.bold = False
        r.font.color.rgb = COLOR_PASS
        p.space_after = Pt(8)

    # Right: Security Audit
    create_card(slide9, 6.8, 1.5, 5.7, 5.2)
    tb9_r = slide9.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.8))
    tf9_r = tb9_r.text_frame
    tf9_r.word_wrap = True

    p = tf9_r.paragraphs[0]
    p.text = "🔒 Security & Isolation Audit"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PASS
    p.space_after = Pt(12)

    sec_items = [
        ("Multi-Tenant Isolation", "Every relationship and responsibility query enforces strict tenant_id filtering. Zero cross-tenant data leaks."),
        ("Zero Hardcoded Secrets", "All DB credentials, JWT secrets, and environment tokens managed strictly via env vars."),
        ("Transactional Integrity", "Single active primary owner rule enforced via atomic DB transactions with rollback on failure."),
        ("Audit Logging Compliance", "All link creations, edits, and revocations publish structured audit events to GovernanceEventService.")
    ]

    for title, desc in sec_items:
        p = tf9_r.add_paragraph()
        p.text = f"• {title}: "
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        r = p.add_run()
        r.text = desc
        r.font.bold = False
        r.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    # ==========================================
    # SLIDE 10: Sign-Off & Deployment Checklist
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10)
    add_header(slide10, "RELEASE APPROVAL & SIGN-OFF", "Production Deployment Readiness Checklist")

    create_card(slide10, 0.8, 1.5, 11.7, 5.2)
    tb10 = slide10.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.1), Inches(4.8))
    tf10 = tb10.text_frame
    tf10.word_wrap = True

    p = tf10.paragraphs[0]
    p.text = "🟢 EXECUTIVE SIGN-OFF: APPROVED FOR PRODUCTION DEPLOYMENT"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PASS
    p.space_after = Pt(16)

    checklist = [
        ("Database Migration", "Run 'alembic upgrade head' to apply consolidated schema migration '7e72dc221571' and composite index migration '73433bbfa6a5'."),
        ("Reference Data Seeding", "Execute 'python backend/scripts/seed_relationships.py' post-migration to seed enterprise entity relationships."),
        ("Governance Monitoring", "Subscribe monitoring alerts to RELATIONSHIP_VALIDATION_FAILED audit events for real-time unauthorized change detection."),
        ("Final Quality Verdict", "22/22 System Features Passed | 6/6 Initial Defect Remediations Verified | 0 Open Bugs Remaining.")
    ]

    for title, desc in checklist:
        p = tf10.add_paragraph()
        p.text = f"✓ {title}: "
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        r = p.add_run()
        r.text = desc
        r.font.bold = False
        r.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(12)

    # Save presentation
    output_path = r"c:\Users\aayus\Desktop\GuardianIQ--1\docs\GuardianIQ_Business_Scenario_QA_Features_Presentation.pptx"
    prs.save(output_path)
    print(f"Successfully generated presentation at: {output_path}")

if __name__ == "__main__":
    create_presentation()
